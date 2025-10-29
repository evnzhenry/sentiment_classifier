"""Generate a stakeholder-friendly PowerPoint that communicates technical concepts
to non-technical audiences while maintaining depth and strategic direction.

Key features:
- High-level objectives and business value first, clear information flow
- Simple analogies, real-world examples, and visual storytelling (minimal text)
- Consistent colors/typography, ample white space, progressive disclosure
- Glossary, FAQ, and clear call-to-action
- Version-controlled speaker notes saved alongside the deck

Run:
    python tools\generate_presentation.py --out AfriSenti_strategy_presentation.pptx --notes speaker_notes.md

Requirements: python-pptx
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import argparse
import json

ROOT = Path('c:/Users/BMC/Desktop/NLP')
DEFAULT_OUT = ROOT / 'AfriSenti_strategy_presentation.pptx'
DEFAULT_NOTES = ROOT / 'speaker_notes.md'


def read_text_file(p: Path):
    if p.exists():
        return p.read_text(encoding='utf-8')
    return ''


def set_theme(slide):
    # Simple theme: white background, dark text, accent color for shapes
    accent = RGBColor(18, 132, 198)  # blue
    for shape in slide.shapes:
        if hasattr(shape, 'fill') and shape.fill:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            except Exception:
                pass
    return accent


def add_title_slide(prs, title, subtitle=''):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    subtitle_tf = slide.placeholders[1]
    title_tf.text = title
    subtitle_tf.text = subtitle
    # Center alignment for readability
    title_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_theme(slide)


def add_text_slide(prs, heading, body, max_chars=600):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    tx = slide.shapes.placeholders[1].text_frame
    # split into paragraphs to keep slides readable
    if len(body) > max_chars:
        body = body[:max_chars] + '\n\n(Truncated; see REPORT.md for full details)'
    for i, line in enumerate(body.split('\n')):
        p = tx.add_paragraph() if i>0 else tx.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
    set_theme(slide)


def add_code_slide(prs, heading, code_snippet, max_lines=20):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    tx = slide.shapes.placeholders[1].text_frame
    lines = code_snippet.split('\n')
    if len(lines) > max_lines:
        lines = lines[:max_lines] + ['# ... (truncated)']
    for i, line in enumerate(lines):
        p = tx.add_paragraph() if i>0 else tx.paragraphs[0]
        p.text = line
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.LEFT
    set_theme(slide)


def add_image_slide(prs, heading, image_path: Path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    try:
        slide.shapes.add_picture(str(image_path), left, top, width=width)
    except Exception as e:
        slide.shapes.placeholders[1].text = f'Image could not be loaded: {e}'
    set_theme(slide)


def gather_eda_images(root: Path):
    candidates = list(root.glob('**/eda*.png')) + list(root.glob('**/char_len_dist.png')) + list(root.glob('**/*_len_dist.png'))
    unique = []
    seen = set()
    for p in candidates:
        if p.name not in seen:
            unique.append(p)
            seen.add(p.name)
    return unique


def add_bullets_slide(prs, heading, bullets, note=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    set_theme(slide)
    if note:
        slide.notes_slide.notes_text_frame.text = note


def add_visual_analogy_slide(prs, title, analogy_lines, diagram_note=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    # Analogy bullets
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
    tf = box.text_frame
    for i, line in enumerate(analogy_lines):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(20)
    # Simple diagram using shapes (infographic style)
    accent = set_theme(slide)
    left = Inches(6.5)
    top = Inches(2)
    width = Inches(3)
    height = Inches(1)
    # Boxes with arrows indicating flow
    box1 = slide.shapes.add_shape(1, left, top, width, height)  # 1=Rectangle
    box1.text_frame.text = 'Data'
    box1.text_frame.paragraphs[0].font.size = Pt(16)
    box1.fill.solid()
    box1.fill.fore_color.rgb = accent
    box1.line.color.rgb = RGBColor(255,255,255)
    box2 = slide.shapes.add_shape(1, left, top+Inches(1.5), width, height)
    box2.text_frame.text = 'Features'
    box2.text_frame.paragraphs[0].font.size = Pt(16)
    box2.fill.solid(); box2.fill.fore_color.rgb = accent
    box3 = slide.shapes.add_shape(1, left, top+Inches(3), width, height)
    box3.text_frame.text = 'Prediction'
    box3.text_frame.paragraphs[0].font.size = Pt(16)
    box3.fill.solid(); box3.fill.fore_color.rgb = accent
    # Notes for deeper technical detail
    if diagram_note:
        slide.notes_slide.notes_text_frame.text = diagram_note


def build_presentation(prs, assets, out_path: Path, notes_path: Path):
    speaker_notes = []
    # Executive Summary
    add_title_slide(prs, 'AfriSenti Strategy Brief', 'Technical impact made accessible')
    add_bullets_slide(
        prs,
        'Executive Summary',
        bullets=[
            'Goal: Reliable multilingual sentiment analysis (Swahili, Amharic, English).',
            'Business value: Faster market insights and customer sentiment across regions.',
            'Approach: Lightweight baselines now, scalable transformers when needed.',
        ],
        note='Summary: Emphasize ROI, coverage, and phased adoption. Provide data volume, label balance, and initial metrics in speaker notes.'
    )
    speaker_notes.append('Exec: Dataset coverage, initial F1/accuracy, and infra assumptions (CPU/GPU).')

    # Objectives & Value
    add_bullets_slide(
        prs,
        'Objectives & Business Value',
        bullets=[
            'Reduce time-to-insight for sentiment shifts across African markets.',
            'Enable cross-language consistency in reporting.',
            'Lay foundation for product feedback loop and brand monitoring.',
        ],
        note='Map objectives to measurable KPIs (dashboard freshness, alert latency, F1 thresholds).'
    )
    speaker_notes.append('KPIs: latency < 1h, macro-F1 ≥ baseline, coverage ≥ target langs.')

    # Concepts via analogies
    add_visual_analogy_slide(
        prs,
        'How It Works (Simple Analogy)',
        ['TF-IDF: like a librarian indexing unique terms (importance).',
         'Logistic Regression: weighs evidence for each sentiment.',
         'SVC: draws clean boundaries between classes.',
         'Transformer models: multilingual understanding out-of-the-box (XLM-R).'],
        diagram_note='Technical: pipeline stages — cleaning, vectorization, training, evaluation; swap in transformer classifier as needed.'
    )
    speaker_notes.append('Detail TF-IDF math, regularization in LogReg, calibration for SVC.')

    # Visual storytelling: EDA
    eda_imgs = assets.get('eda_imgs', [])
    if eda_imgs:
        add_bullets_slide(prs, 'Data Story (EDA)', ['Balanced labels drive reliable metrics.', 'Language mix informs reporting strategy.', 'Text lengths guide preprocessing.'], note='Refer to EDA plots on next slides.')
        for img in eda_imgs[:4]:
            add_image_slide(prs, f'EDA Visual: {img.stem}', img)
        speaker_notes.append('EDA details: label distribution, per-language counts, char/token stats.')
    else:
        add_text_slide(prs, 'Data Story (EDA)', 'Run EDA to populate visuals: label/lang distributions and text length histograms.')
        speaker_notes.append('EDA pending; provide dataset stats when available.')

    # PCA simplified
    pca_path = assets.get('pca_img')
    if pca_path:
        add_image_slide(prs, 'Feature Space Snapshot (PCA)', pca_path)
        speaker_notes.append('PCA EVR and top component loadings in real_pca_evr.json and real_pca_loadings.json.')
    else:
        add_text_slide(prs, 'Feature Space Snapshot (PCA)', 'PCA gives a bird’s-eye view of how sentiments cluster. Run PCA to visualize.')
        speaker_notes.append('PCA pending; include variance and loadings when generated.')

    # Baseline results, minimal text
    metrics_summary = assets.get('metrics_summary')
    if metrics_summary:
        add_bullets_slide(prs, 'Baseline Performance (Simple, Strong Start)', [
            f"LogReg macro F1: {metrics_summary.get('tfidf_logreg_macro_f1','N/A')}",
            f"Naive Bayes macro F1: {metrics_summary.get('tfidf_nb_macro_f1','N/A')}",
            f"Calibrated SVC macro F1: {metrics_summary.get('tfidf_svc_macro_f1','N/A')}",
        ], note='Use per-class F1 bar chart for category-level clarity (model_f1_bar_real.png).')
        speaker_notes.append('Include ROC-AUC if available; discuss trade-offs and deployment fit.')
        if assets.get('f1_bar'):
            add_image_slide(prs, 'Per-class F1 by Model', assets['f1_bar'])
    else:
        add_text_slide(prs, 'Baseline Performance', 'Run the real-data pipeline to populate macro F1 and confusion matrices.')
        speaker_notes.append('Awaiting metrics; once ready, add thresholds for green/yellow/red statuses.')

    # Strategic Direction & Roadmap
    add_bullets_slide(prs, 'Strategic Direction', [
        'Phase 1: Baselines for quick wins and reporting.',
        'Phase 2: Fine-tune multilingual transformers for improved accuracy.',
        'Phase 3: Production monitoring, drift detection, retraining cadence.',
    ], note='Detail data refresh, feedback loops, and incident response.')
    speaker_notes.append('Infra: batch inference on CPU, transformer finetuning on GPU when needed.')

    add_bullets_slide(prs, 'Implementation Roadmap', [
        'Week 1–2: EDA, baselines, dashboards.',
        'Week 3–4: Model hardening and transformer experiments.',
        'Week 5+: Deployment, observability, and stakeholder training.',
    ], note='Deliverables: reports, visuals, KPI dashboard, training materials.')
    speaker_notes.append('Owners, timelines, and dependency tracking.')

    # Glossary
    add_bullets_slide(prs, 'Glossary', [
        'TF-IDF: Importance score for words based on rarity and usage.',
        'Macro F1: Balanced average of F1 across labels.',
        'ROC-AUC: Probability a model ranks positives over negatives.',
        'Transformer: Neural architecture with attention enabling language understanding.',
    ], note='Include links or references in REPORT.md for further reading.')
    speaker_notes.append('Provide brief math references and citations.')

    # FAQ
    add_bullets_slide(prs, 'FAQ', [
        'How accurate is it? Baselines provide reliable starting points.',
        'What data do we need? Labeled examples per language (balanced).',
        'How do we deploy? Batch scoring first; online scoring later if needed.',
        'What about bias? Monitor per-class/per-language performance; retrain as needed.',
    ], note='Address privacy, compliance, and language evolution.')
    speaker_notes.append('Bias and fairness checks, data governance, opt-in/opt-out policies.')

    # Call to Action
    add_bullets_slide(prs, 'Call to Action', [
        'Approve Phase 1 baseline deployment to insights dashboard.',
        'Schedule stakeholder review of KPIs and roadmap.',
        'Greenlight data refresh cadence and model monitoring.',
    ], note='Close with business stakeholders’ commitments and timelines.')
    speaker_notes.append('Define decision gates, budget, and resource allocations.')

    # Attach minimal code slide for technical audiences (optional)
    preprocess_txt = assets.get('preprocess_txt')
    if preprocess_txt:
        add_code_slide(prs, 'Preprocessing (for engineers)', '\n'.join(preprocess_txt.splitlines()[:50]))
        speaker_notes.append('Detail cleaning steps and configuration options.')

    # Save deck
    prs.save(str(out_path))
    # Save speaker notes as markdown for version control
    notes_path.write_text('\n'.join(f"- {n}" for n in speaker_notes), encoding='utf-8')
    print(f'Presentation saved to: {out_path}')
    print(f'Speaker notes saved to: {notes_path}')


def main():
    parser = argparse.ArgumentParser(description='Generate stakeholder-friendly AfriSenti presentation')
    parser.add_argument('--out', type=str, default=str(DEFAULT_OUT), help='Output PPTX path')
    parser.add_argument('--notes', type=str, default=str(DEFAULT_NOTES), help='Speaker notes markdown path')
    args = parser.parse_args()

    assets = {}
    # Gather images and artifacts
    eda_imgs = gather_eda_images(ROOT / 'eda_outputs')
    assets['eda_imgs'] = eda_imgs
    pca_img = (ROOT / 'eda_outputs' / 'real_pca_tfidf_scatter.png')
    assets['pca_img'] = pca_img if pca_img.exists() else None
    f1_bar = (ROOT / 'eda_outputs' / 'model_f1_bar_real.png')
    assets['f1_bar'] = f1_bar if f1_bar.exists() else None
    # Metrics summary (if available)
    metrics = {}
    for name in ['tfidf_logreg', 'tfidf_nb', 'tfidf_svc']:
        p = ROOT / 'eda_outputs' / f'{name}_metrics_real.json'
        if p.exists():
            try:
                data = json.loads(p.read_text(encoding='utf-8'))
                metrics[f'{name}_macro_f1'] = data.get('macro_f1', 'N/A')
            except Exception:
                metrics[f'{name}_macro_f1'] = 'N/A'
    assets['metrics_summary'] = metrics if metrics else None
    assets['preprocess_txt'] = read_text_file(ROOT / 'src' / 'preprocess.py')

    prs = Presentation()
    build_presentation(prs, assets, Path(args.out), Path(args.notes))


if __name__ == '__main__':
    main()
